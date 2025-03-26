import os
import pandas as pd
import cv2
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation

# Set the path to Tesseract executable
# Windows example:
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
# Linux/Mac example:
# pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

def extract_text_from_file(file_path):
    """Extracts text from PDF, Images, Excel, and PowerPoint files."""
    ext = file_path.lower().split(".")[-1]
    extracted_text = ""

    if ext in ["png", "jpg", "jpeg"]:
        print("$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$")
        try:
            # Image OCR using pytesseract
            image = cv2.imread(file_path)
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            gray = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]

            # Use pytesseract for text extraction
            extracted_text = pytesseract.image_to_string(gray)

            print("$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$")
            print(extracted_text)

        except Exception as e:
            print(f"Error processing image: {e}")
            extracted_text = ""

    elif ext == "pdf":
        try:
            with pdfplumber.open(file_path) as pdf:
                extracted_text = "\n".join(
                    page.extract_text() for page in pdf.pages if page.extract_text()
                )
        except Exception as e:
            print(f"Error processing PDF: {e}")
            extracted_text = ""

    elif ext in ["xls", "xlsx"]:
        try:
            df = pd.read_excel(file_path, sheet_name=None)
            extracted_text = "\n".join(
                df[sheet].to_string() for sheet in df
            )
        except Exception as e:
            print(f"Error processing Excel: {e}")
            extracted_text = ""

    elif ext == "pptx":
        try:
            presentation = Presentation(file_path)
            extracted_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        extracted_text.append(shape.text.strip())
            extracted_text = "\n".join(extracted_text)
        except Exception as e:
            print(f"Error processing PowerPoint file: {e}")
            extracted_text = ""

    return extracted_text.strip()

